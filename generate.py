from bullet import Bullet, Check, YesNo, Input, VerticalPrompt # and etc...
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import PP_ALIGN
import sys

TITLE_LAYOUT = 0
BLANK = 10
FONT = "Calibri"
FILE_LOC = ""
SERMON_SERIES = "Signs of Life Series in John’s Gospel"

cli = VerticalPrompt(
    [
        Input("What is the reference for Call to Worship?"),
        Input("Filename?"),
        Input("What is the reference for Confession of Sin?"),
        Input("Filename?"),
        Input("What is the reference for Assurance of Worship?"),
        Input("Filename?"),
        YesNo("Is it prayers of the people? "),
        Input("Sermon title?"),
        Input("Sermon reference?"),
        #Bullet("What is your favorite programming language? ",
         #     choices = ["C++", "Python", "Javascript", "Not here!"]),
    ],
    spacing = 1
)

result = cli.launch()

class LiturgyElem:
    def __init__(self,f_name):
        self.title = ""
        self.lyrics = []
        self.f_name = f_name

    #opens a file and turns it into a LiturgyElem
    def parseLiturgyFile(self,delimiter):
        with open(self.f_name,"r") as file:
            lyrics = file.read()
            lyrics = lyrics.split(delimiter)
            self.title = lyrics[0]
            self.lyrics = lyrics[1:]

def parseFileLines(filename):
    with open(filename,"r") as file:
        return file.readlines()

def parseFile(filename):
    with open(filename,"r") as file:
        return file.read()

#creates a textbox for a song slide, with parameters for placement, size, and boldness
#returns a paragraph!!
def addTextFrame(slide,dim):
        shapes = slide.shapes
        verse_text = shapes.add_textbox(Inches(dim[0]),Inches(dim[1]),Inches(dim[2]),Inches(dim[3]))
        frame = verse_text.text_frame
        frame.word_wrap = True
        return frame

#adds a center justified paragraph with given font size and boldness to the given frame
def addCenterParagraph(frame,size,bold):
    if (frame.paragraphs[0].font.name == 'Calibri'):
        p = frame.add_paragraph()
    else:
        p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Calibri'
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(255,255,255)
    return p

#adds a left justified paragraph with given font size and boldness to the given frame
def addLeftParagraph(frame,size,bold):
    if (frame.paragraphs[0].font.name == 'Calibri'):
        p = frame.add_paragraph()
    else:
        p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.name = 'Calibri'
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(255,255,255)
    return frame

#adds a song to the presentation at the end, pulling from a LiturgyElem object
def addSong(prs,filename):
    #add title slide for song
    song = LiturgyElem(filename)
    song.parseLiturgyFile("\n\n")
    song_title_slide = addBlank(prs)
    title = addTextFrame(song_title_slide,(0.375,0.25,9.25,1.25))
    p = addCenterParagraph(title,36,True)
    title.text = song.title
    #put first verse of song on title slide
    frame = addTextFrame(song_title_slide,(0.5,1.5,9,6))
    p = addCenterParagraph(frame,30,False)
    p.text = song.lyrics[0]
    #loop through remaining verses and add slides
    for verse in song.lyrics[1:]:
        verse_slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
        frame = addTextFrame(verse_slide,(0.5,1.5,9,6))
        p = addCenterParagraph(frame,30,False)
        p.text = verse



def addResponsiveSlide(prs,grace,segment):
    i = 0
    #add title slide for segment
    first_slide = addBlank(prs)
    frame = addTextFrame(first_slide,(0.5,0.25,9,7.25))
    title = addCenterParagraph(frame,30,True)
    #assurance,confession,call
    title.text = segment
    if grace:
        mid = addCenterParagraph(frame,24,False)
        mid.text = "Now hear the gracious word of God from:"
    sub = addCenterParagraph(frame,30,True)
    sub.text = result[i][1]
    i+=1
    filename = result[i][1]
    sections = parseFile(filename).split("$$$")
    #put first section of response
    frame = addTextFrame(first_slide,(0.5,1.5,9,6))
    p = addLeftParagraph(frame,28,False)
    p.text = sections[0]
    #loop through remaining sections and add slides
    for section in sections[1:]:
        slide = addBlank(prs)
        frame = addTextFrame(slide,(0.5,0.75,9,6.75))
        p = addLeftParagraph(frame,28,False)
        p.text = section
    i+=1


def addBlank(prs):
    return prs.slides.add_slide(prs.slide_layouts[BLANK])

#open template
prs = Presentation("template.pptx")
addBlank(prs)

#welcome slide
welcome_slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
title_placeholder = welcome_slide.placeholders[0]
title_placeholder.text = "Welcome to Living Hope Church!"
addBlank(prs)

#first song
filename = FILE_LOC + input("Name of first song? ")
addSong(prs,filename)
addBlank(prs)

#call to worship
addResponsiveSlide(prs,False,"Call to Worship")
addBlank(prs)

#second, third, fourth song
for i in range(2):
    filename = FILE_LOC + input("Name of song? ")
    addSong(prs,filename)
    addBlank(prs)

#confession
addResponsiveSlide(prs,False,"Confession of Sin")
addBlank(prs)

#assurance of grace
addResponsiveSlide(prs,False,"Assurance of Grace")
addBlank(prs)

#offering
offering_slide = addBlank(prs)
frame = addTextFrame(offering_slide,(0,0.5,10,1.5))
p = addCenterParagraph(frame,32,True)
p.text = "Worship through Tithes and Offerings"
frame = addTextFrame(offering_slide,(0,1.5,10,6))
p = addCenterParagraph(frame,28,False)
p.text = parseFile(FILE_LOC + "Offering.txt")
addBlank(prs)

#offering song
filename = FILE_LOC + input("Name of song? ")
addSong(prs,filename)
addBlank(prs)

#Prayers of the people
prayer_slide = addBlank(prs)
frame = addTextFrame(prayer_slide,(0,1.5,10,1.5))
p = addCenterParagraph(frame,30,True)
if (result[6][1]):
    p.text = "Prayers of the People"
else:
    p.text = "Prayers for our Church, Community, and World"
addBlank(prs)

#Announcements
ann_slide = addBlank(prs)
frame = addTextFrame(ann_slide,(0.25,0,9,7))
p = addCenterParagraph(frame,28,True)
p.text = "Announcements"
addBlank(prs)

#Sermon
sermon_slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
frame = addTextFrame(sermon_slide,(0.75,1.5,8.5,4.5))
title = addCenterParagraph(frame,30,True)
title.text = result[7][1]
mid = addCenterParagraph(frame,24,True)
mid.font.input = True
mid.text = SERMON_SERIES
sub = addCenterParagraph(frame,30,True)
sub.text = result[8][1]

#closing song
filename = FILE_LOC + input("Name of song? ")
addSong(prs,filename)
addBlank(prs)

bene_slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
frame = addTextFrame(bene_slide,(0,0.5,10,1))
p = addCenterParagraph(frame,44,True)
p.text = "Benediction"
rest = addTextFrame(bene_slide,(0,1.75,10,3.5))
p = addCenterParagraph(rest,40,False)
p.text = "A blessing from God’s Word\nas we go out into God’s world."

prs.save("slidesbase.pptx")
